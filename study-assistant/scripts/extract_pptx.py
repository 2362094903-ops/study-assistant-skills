#!/usr/bin/env python3
"""Extract text content from PowerPoint (.pptx) files as Markdown.

Why: teachers often distribute courseware as PPT slides. This script extracts
slide titles, body text, tables, and speaker notes into the same Markdown format
as extract_pdf.py, so the orchestrator can treat them interchangeably.
Image-heavy slides (text < 50 chars + has images) are flagged for OCR.

Usage:
  python3 extract_pptx.py courseware.pptx -o chapter-03.md
  python3 extract_pptx.py courseware.pptx --slides 1-15 -o chapter-03.md
  python3 extract_pptx.py courseware.pptx --slides 1-15 -o ch3.md --render-images ./slides/
"""

import argparse
import pathlib
import re
import sys

IMAGE_HEAVY_THRESHOLD = 50  # same as extract_pdf.py SCANNED_THRESHOLD


# ---------------------------------------------------------------------------
# Slide range parsing (identical to extract_pdf.parse_pages)
# ---------------------------------------------------------------------------
def parse_slides(spec, total):
    """Parse a 1-indexed page/slide range spec.  Examples: "1-5", "3", "1-5,8,10-12"."""
    if not spec:
        return list(range(total))
    result = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            result.extend(range(int(a.strip()) - 1, int(b.strip())))
        else:
            result.append(int(part) - 1)
    return [i for i in result if 0 <= i < total]


# ---------------------------------------------------------------------------
# Slide content extraction
# ---------------------------------------------------------------------------
def extract_slide(slide, img_dir=None, slide_idx=0, img_counter=None):
    """Extract text and metadata from a single slide.

    Returns a dict:
      {"title": str, "body": [str], "tables": [[[str]]], "notes": str,
       "image_count": int, "char_count": int}
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Inches
    except ImportError:
        sys.exit("缺少 python-pptx，请先运行: pip3 install python-pptx")

    result = {
        "title": "",
        "body": [],
        "tables": [],
        "notes": "",
        "image_count": 0,
        "char_count": 0,
    }

    # Check for title placeholder
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                t = shape.text_frame.text.strip() if shape.has_text_frame else ""
                if t:
                    result["title"] = t
                break
    except Exception:
        pass

    for shape in slide.shapes:
        # --- text shapes ---
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            char_count = len(text)
            result["char_count"] += char_count

            # Skip if this was already captured as title
            if text and text == result["title"]:
                continue

            # If title is empty and this looks like a title, use it
            if not result["title"] and char_count < 120 and char_count > 0:
                is_title_candidate = True
                try:
                    from pptx.enum.shapes import PP_PLACEHOLDER
                    if shape.is_placeholder:
                        if shape.placeholder_format.type not in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                            is_title_candidate = False
                except Exception:
                    pass
                if is_title_candidate:
                    result["title"] = text
                    continue

            if text:
                result["body"].append(text)

        # --- tables ---
        if shape.has_table:
            table = shape.table
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip().replace("\n", " ")
                    result["char_count"] += len(cell_text)
                    cells.append(cell_text)
                rows.append(cells)
            if rows:
                result["tables"].append(rows)

        # --- images ---
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                result["image_count"] += 1
                # Export image if requested
                if img_dir and img_counter is not None:
                    try:
                        img = shape.image
                        ext = img.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        fname = f"slide-{slide_idx + 1}-img-{img_counter[0]}.{ext}"
                        (img_dir / fname).write_bytes(img.blob)
                        img_counter[0] += 1
                    except Exception:
                        pass
        except Exception:
            pass

        # --- grouped shapes: recurse ---
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.has_text_frame:
                        text = child.text_frame.text.strip()
                        if text and text != result["title"]:
                            result["body"].append(text)
                            result["char_count"] += len(text)
        except Exception:
            pass

    # --- speaker notes ---
    try:
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                result["notes"] = notes_text
                result["char_count"] += len(notes_text)
    except Exception:
        pass

    return result


# ---------------------------------------------------------------------------
# Markdown rendering
# ---------------------------------------------------------------------------
def render_markdown_table(rows):
    """Render a 2D list as a pipe-separated Markdown table."""
    if not rows:
        return []
    lines = []
    # Header row
    lines.append("| " + " | ".join(rows[0]) + " |")
    # Separator
    lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    # Data rows
    for row in rows[1:]:
        # Pad shorter rows
        padded = row + [""] * (len(rows[0]) - len(row))
        lines.append("| " + " | ".join(padded) + " |")
    return lines


def slide_to_markdown(data, slide_num, img_dir=None):
    """Convert extracted slide data to Markdown lines."""
    lines = [f"<!-- slide {slide_num} -->"]

    # Image-heavy detection
    if data["char_count"] < IMAGE_HEAVY_THRESHOLD and data["image_count"] > 0:
        lines.append(f"<!-- slide {slide_num}: IMAGE-HEAVY, 待 OCR -->")

    if data["title"]:
        lines.append(f"# {data['title']}")
        lines.append("")

    for para in data["body"]:
        lines.append(para)
        lines.append("")

    for table in data["tables"]:
        lines.extend(render_markdown_table(table))
        lines.append("")

    if data["notes"]:
        lines.append(f"> **讲稿备注：** {data['notes']}")
        lines.append("")

    return lines


# ---------------------------------------------------------------------------
# .ppt fallback via LibreOffice
# ---------------------------------------------------------------------------
def convert_ppt_via_libreoffice(ppt_path):
    """Try to convert .ppt to .pptx using LibreOffice. Returns Path to .pptx or None."""
    import subprocess
    import tempfile
    import shutil

    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice:
        return None

    out_dir = pathlib.Path(tempfile.mkdtemp())
    try:
        result = subprocess.run(
            [libreoffice, "--headless", "--convert-to", "pptx",
             "--outdir", str(out_dir), str(ppt_path)],
            capture_output=True, text=True, timeout=60)
        converted = list(out_dir.glob("*.pptx"))
        if converted:
            return converted[0]
    except Exception:
        pass
    return None


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser(
        description="Extract PPTX slide text as Markdown (with slide markers)")
    ap.add_argument("pptx_file", help="path to .pptx (or .ppt) file")
    ap.add_argument("--slides", help="slide range, 1-indexed, e.g. 1-15 or 1-5,8")
    ap.add_argument("-o", "--out", required=True, help="output markdown path")
    ap.add_argument("--render-images", metavar="DIR",
                    help="export embedded slide images to this directory for OCR")
    args = ap.parse_args()

    pptx_path = pathlib.Path(args.pptx_file)
    if not pptx_path.exists():
        sys.exit(f"文件不存在：{pptx_path}")

    # .ppt fallback
    is_old = pptx_path.suffix.lower() == ".ppt"
    if is_old:
        print("⚠️  .ppt 是旧版格式，python-pptx 不支持，尝试用 LibreOffice 转换...",
              file=sys.stderr)
        converted = convert_ppt_via_libreoffice(pptx_path)
        if converted:
            pptx_path = converted
            print(f"   已转换为 .pptx：{pptx_path}", file=sys.stderr)
        else:
            sys.exit(
                "无法处理 .ppt 文件。\n"
                "  方案 1：用 PowerPoint / WPS 另存为 .pptx 格式\n"
                "  方案 2：安装 LibreOffice 自动转换：brew install libreoffice")

    # Import python-pptx
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("缺少 python-pptx，请先运行: pip3 install python-pptx")

    prs = Presentation(str(pptx_path))
    total = len(prs.slides)
    slide_indices = parse_slides(args.slides, total)

    img_dir = None
    img_counter = None
    if args.render_images:
        img_dir = pathlib.Path(args.render_images)
        img_dir.mkdir(parents=True, exist_ok=True)
        img_counter = [1]

    image_heavy_slides = []
    extracted_count = 0
    lines = []

    for idx in slide_indices:
        slide = prs.slides[idx]
        data = extract_slide(slide, img_dir, idx, img_counter)
        slide_num = idx + 1

        slide_lines = slide_to_markdown(data, slide_num, img_dir)
        lines.extend(slide_lines)

        if data["char_count"] < IMAGE_HEAVY_THRESHOLD and data["image_count"] > 0:
            image_heavy_slides.append(slide_num)

        extracted_count += 1

    # Write output
    out = pathlib.Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text("\n".join(lines), encoding="utf-8")

    # Summary to stderr
    print(f"已提取 {extracted_count} 页幻灯片 → {out}", file=sys.stderr)
    if image_heavy_slides:
        where = f"，图片已导出到 {args.render_images}" if args.render_images else ""
        print(f"疑似图片为主幻灯片（需 OCR）：{image_heavy_slides}{where}",
              file=sys.stderr)
    else:
        print("未发现纯图片幻灯片。", file=sys.stderr)


if __name__ == "__main__":
    main()
